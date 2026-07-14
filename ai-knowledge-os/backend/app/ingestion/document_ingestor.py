import os
import uuid
from typing import List, Dict, Any, Optional
from bs4 import BeautifulSoup
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from sentence_transformers import SentenceTransformer

from app.services.vector_store import VectorStoreService

from app.config import settings

class DocumentIngestor:
    def __init__(self, vector_store: Optional[VectorStoreService] = None, model_name: Optional[str] = None):
        """
        Initializes the Document Ingestor with a vector store service and a local embedding model.
        """
        self.vector_store = vector_store or VectorStoreService()
        # SentenceTransformer loads the model locally (downloads on first run)
        model = model_name or settings.EMBEDDING_MODEL_NAME
        self.model = SentenceTransformer(model)
        self.vector_size = self.model.get_embedding_dimension()
        
    def extract_text(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from a file, returning a list of pages/sections with metadata.
        Output format:
        [
            {
                "text": "page or section text...",
                "metadata": {"page_number": int, "source": str}
            }
        ]
        """
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        sections = []
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
            
        if ext == ".pdf":
            doc = fitz.open(file_path)
            for i, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    sections.append({
                        "text": text,
                        "metadata": {"page_number": i + 1, "source": filename}
                    })
            doc.close()
            
        elif ext == ".docx":
            doc = docx.Document(file_path)
            # Group text into paragraphs as sections
            paragraph_texts = [p.text for p in doc.paragraphs if p.text.strip()]
            for idx, text in enumerate(paragraph_texts):
                sections.append({
                    "text": text,
                    "metadata": {"section_index": idx + 1, "source": filename}
                })
                
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for idx, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                combined_text = "\n".join(slide_text)
                if combined_text.strip():
                    sections.append({
                        "text": combined_text,
                        "metadata": {"slide_number": idx + 1, "source": filename}
                    })
                    
        elif ext in [".html", ".htm"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                text = soup.get_text()
                # Clean up whitespace
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                clean_text = "\n".join(chunk for chunk in chunks if chunk)
                if clean_text.strip():
                    sections.append({
                        "text": clean_text,
                        "metadata": {"source": filename}
                    })
                    
        else: # Default/fallback for text files
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
                if text.strip():
                    sections.append({
                        "text": text,
                        "metadata": {"source": filename}
                    })
                    
        return sections

    def chunk_text(self, text: str, chunk_size: int = 500, chunk_overlap: int = 100) -> List[str]:
        """
        Splits a string of text into smaller overlapping chunks.
        """
        if not text:
            return []
            
        chunks = []
        start = 0
        text_len = len(text)
        
        while start < text_len:
            # Take a chunk of chunk_size characters
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            
            # If the current chunk reached or exceeded the end, we are done
            if end >= text_len:
                break
                
            # Move the start index forward by chunk_size - chunk_overlap
            start += (chunk_size - chunk_overlap)
            if start >= text_len or chunk_size <= chunk_overlap:
                break
                
        return chunks

    def ingest_file(self, file_path: str, collection_name: str, chunk_size: int = 500, chunk_overlap: int = 100) -> Dict[str, Any]:
        """
        Extracts, chunks, embeds, and stores a document file into the vector store.
        """
        filename = os.path.basename(file_path)
        
        # 1. Extract text from file
        sections = self.extract_text(file_path)
        if not sections:
            return {"status": "skipped", "message": "No text extracted from file", "filename": filename}
            
        # Ensure collection exists in Qdrant
        self.vector_store.create_collection(collection_name, vector_size=self.vector_size)
        
        all_points = []
        total_chunks = 0
        
        # 2. Chunk text and prepare points
        for sec in sections:
            sec_text = sec["text"]
            sec_metadata = sec["metadata"]
            
            chunks = self.chunk_text(sec_text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
            
            for chunk_idx, chunk in enumerate(chunks):
                # Generate a unique ID for the vector point
                point_id = str(uuid.uuid4())
                
                # Combine section metadata with chunk metadata
                payload = {
                    **sec_metadata,
                    "text": chunk,
                    "chunk_index": chunk_idx,
                    "length": len(chunk)
                }
                
                # Generate embedding
                embedding = self.model.encode(chunk).tolist()
                
                all_points.append({
                    "id": point_id,
                    "vector": embedding,
                    "payload": payload
                })
                
                total_chunks += 1
                
        # 3. Upsert to Qdrant
        if all_points:
            success = self.vector_store.upsert_vectors(collection_name, all_points)
            if success:
                return {
                    "status": "success",
                    "filename": filename,
                    "collection": collection_name,
                    "total_chunks": total_chunks
                }
                
        return {"status": "error", "message": "Failed to store vectors", "filename": filename}
